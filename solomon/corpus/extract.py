"""File-type-specific text extraction.

Ported from /root/projects/solomon-from-drive/corpus_ingest/extract.py per
REPORT-CORPUS.md §1.2 and §4. Each extractor lazy-imports its dependency
so the module loads even when an optional package is missing; the missing
dep raises ``UnsupportedFileType`` which the orchestrator handles by parking
the file under ``corpus/inbox/_unsupported/``.

Public surface:
  - ``extract(path) -> ExtractedDoc``  — dispatch + metadata wrapper
  - ``extract_pdf / extract_docx / ...`` — individual extractors
  - ``UnsupportedFileType``            — raised on missing deps / bad files

Sonnet multimodal fallback for scanned PDFs / images is documented in
REPORT-CORPUS.md §4.6 as the lone remote dependency. We stub it behind a
feature flag (``SOLOMON_ALLOW_VISION_API``); when disabled, scanned
PDFs / images raise ``UnsupportedFileType`` and the orchestrator parks
them.
"""

from __future__ import annotations

import csv
import json
import logging
import os
import re
from dataclasses import dataclass, field
from email import policy
from email.parser import BytesParser
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger("solomon.corpus.extract")


# ---------------------------------------------------------------------------
# Exceptions and dataclass
# ---------------------------------------------------------------------------


class UnsupportedFileType(Exception):
    """Raised when a file can't be turned into text.

    Caller (the ingest orchestrator) parks the offending file under
    ``corpus/inbox/_unsupported/`` with the message attached, rather than
    crashing the whole worker.
    """


@dataclass
class ExtractedDoc:
    """Result of extract(path).

    Attributes
    ----------
    text:
        The full extracted unicode text. Always non-empty (we raise
        UnsupportedFileType if extraction would yield an empty string).
    page_breaks:
        Sorted list of character offsets where a new logical "page"
        (PDF page, PPTX slide, XLSX sheet, EML part) begins. Empty for
        flat-text formats.
    metadata:
        Free-form dict. Keys we set today:
        ``extractor`` (one of 'pdf', 'docx', ...), ``format`` (original
        suffix), ``page_count`` (where applicable), ``via`` (set to
        'vision' when we fell back to multimodal — currently always
        absent because the stub never runs).
    """

    text: str
    page_breaks: List[int] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)


# ---------------------------------------------------------------------------
# Extension buckets — kept identical to the Drive version
# ---------------------------------------------------------------------------

PLAIN_TEXT_EXT = {".txt", ".md"}
RTF_EXT = {".rtf"}
PDF_EXT = {".pdf"}
DOCX_EXT = {".docx"}
PPTX_EXT = {".pptx"}
XLSX_EXT = {".xlsx"}
HTML_EXT = {".html", ".htm"}
EMAIL_EXT = {".eml", ".mbox"}
CSV_EXT = {".csv", ".tsv"}
JSON_EXT = {".json"}
IMAGE_EXT = {".png", ".jpg", ".jpeg", ".heic", ".webp", ".gif"}


# ---------------------------------------------------------------------------
# Plain text + Markdown
# ---------------------------------------------------------------------------


def extract_txt_md(path: Path) -> ExtractedDoc:
    """Read a UTF-8 (lossy) text or markdown file."""
    text = path.read_text(encoding="utf-8", errors="replace")
    return ExtractedDoc(
        text=text,
        metadata={"extractor": "txt_md", "format": path.suffix.lower()},
    )


# ---------------------------------------------------------------------------
# RTF — best-effort via striprtf if installed, else return as plain text
# ---------------------------------------------------------------------------


def extract_rtf(path: Path) -> ExtractedDoc:
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        from striprtf.striprtf import rtf_to_text  # type: ignore
        text = rtf_to_text(raw)
    except ImportError:
        logger.warning("striprtf not installed; reading .rtf as raw text")
        text = raw
    return ExtractedDoc(text=text, metadata={"extractor": "rtf"})


# ---------------------------------------------------------------------------
# PDF (pypdf)
# ---------------------------------------------------------------------------


def extract_pdf(path: Path) -> ExtractedDoc:
    """Try the text layer first; fall back to vision only when allowed."""
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise UnsupportedFileType("pypdf not installed") from e

    reader = PdfReader(str(path))
    pages: List[str] = []
    for page in reader.pages:
        try:
            pages.append(page.extract_text() or "")
        except Exception:  # noqa: BLE001
            logger.exception("PDF page extraction failed; continuing")
            pages.append("")

    # Compute page-break offsets while joining.
    parts: List[str] = []
    breaks: List[int] = []
    cursor = 0
    for i, p in enumerate(pages):
        p_stripped = p.strip()
        if not p_stripped:
            continue
        if parts:
            sep = "\n\n"
            cursor += len(sep)
            parts.append(sep)
        breaks.append(cursor)
        cursor += len(p_stripped)
        parts.append(p_stripped)
    text = "".join(parts)

    if text.strip():
        return ExtractedDoc(
            text=text,
            page_breaks=breaks,
            metadata={
                "extractor": "pdf",
                "page_count": len(pages),
            },
        )

    # No text layer — scanned PDF. Fall back to Sonnet multimodal stub.
    logger.info("PDF has no text layer; routing to multimodal fallback")
    try:
        vtext = extract_text_via_sonnet(path)
    except UnsupportedFileType:
        raise
    except Exception as e:  # noqa: BLE001
        raise UnsupportedFileType(f"scanned PDF, vision fallback failed: {e}") from e
    return ExtractedDoc(
        text=vtext,
        metadata={"extractor": "pdf", "page_count": len(pages), "via": "vision"},
    )


# ---------------------------------------------------------------------------
# DOCX (python-docx)
# ---------------------------------------------------------------------------


def extract_docx(path: Path) -> ExtractedDoc:
    try:
        import docx  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("python-docx not installed") from e

    doc = docx.Document(str(path))
    parts: List[str] = []
    for p in doc.paragraphs:
        if p.text and p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    text = "\n".join(parts)
    return ExtractedDoc(text=text, metadata={"extractor": "docx"})


# ---------------------------------------------------------------------------
# PPTX (python-pptx) — one logical "page" per slide
# ---------------------------------------------------------------------------


def extract_pptx(path: Path) -> ExtractedDoc:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("python-pptx not installed") from e

    pres = Presentation(str(path))
    parts: List[str] = []
    breaks: List[int] = []
    cursor = 0
    slide_count = 0
    for i, slide in enumerate(pres.slides, start=1):
        slide_count = i
        header = f"## Slide {i}"
        if parts:
            sep = "\n\n"
            cursor += len(sep)
            parts.append(sep)
        breaks.append(cursor)
        cursor += len(header)
        parts.append(header)
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                parts.append("\n")
                cursor += 1
                parts.append(text)
                cursor += len(text)
    return ExtractedDoc(
        text="".join(parts),
        page_breaks=breaks,
        metadata={"extractor": "pptx", "page_count": slide_count},
    )


# ---------------------------------------------------------------------------
# XLSX (openpyxl) — one "page" per sheet
# ---------------------------------------------------------------------------


def extract_xlsx(path: Path) -> ExtractedDoc:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError as e:
        raise UnsupportedFileType("openpyxl not installed") from e

    wb = load_workbook(filename=str(path), data_only=True, read_only=True)
    parts: List[str] = []
    breaks: List[int] = []
    cursor = 0
    for sheet in wb.sheetnames:
        header = f"## Sheet: {sheet}"
        if parts:
            sep = "\n\n"
            cursor += len(sep)
            parts.append(sep)
        breaks.append(cursor)
        cursor += len(header)
        parts.append(header)
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                line = "\n" + " | ".join(cells)
                parts.append(line)
                cursor += len(line)
    return ExtractedDoc(
        text="".join(parts),
        page_breaks=breaks,
        metadata={"extractor": "xlsx", "page_count": len(wb.sheetnames)},
    )


# ---------------------------------------------------------------------------
# HTML — beautifulsoup4 (preferred) or stdlib HTMLParser fallback
# ---------------------------------------------------------------------------


class _HTMLStripper(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: List[str] = []
        self._skip = False

    def handle_starttag(self, tag: str, attrs):  # type: ignore[override]
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag: str):  # type: ignore[override]
        if tag in ("script", "style"):
            self._skip = False

    def handle_data(self, data: str):  # type: ignore[override]
        if not self._skip and data.strip():
            self.parts.append(data.strip())


def extract_html(path: Path) -> ExtractedDoc:
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        from bs4 import BeautifulSoup  # type: ignore
    except ImportError:
        p = _HTMLStripper()
        p.feed(raw)
        return ExtractedDoc(text="\n".join(p.parts), metadata={"extractor": "html"})
    soup = BeautifulSoup(raw, "html.parser")
    # Drop noise tags entirely.
    for t in soup(["script", "style", "noscript"]):
        t.decompose()
    text = soup.get_text(separator="\n").strip()
    # Collapse runs of blank lines.
    text = re.sub(r"\n{3,}", "\n\n", text)
    return ExtractedDoc(text=text, metadata={"extractor": "html"})


# ---------------------------------------------------------------------------
# Email (.eml / .mbox)
# ---------------------------------------------------------------------------


def extract_eml(path: Path) -> ExtractedDoc:
    raw = path.read_bytes()
    msg = BytesParser(policy=policy.default).parsebytes(raw)
    header_lines: List[str] = []
    headers_meta: Dict[str, Optional[str]] = {}
    for key in ("From", "To", "Cc", "Subject", "Date"):
        v = msg.get(key)
        if v:
            header_lines.append(f"{key}: {v}")
            headers_meta[key.lower()] = str(v)

    body = ""
    if msg.is_multipart():
        # Prefer text/plain
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                body = part.get_content()
                break
        if not body:
            for part in msg.walk():
                if part.get_content_type() == "text/html":
                    stripper = _HTMLStripper()
                    stripper.feed(part.get_content())
                    body = "\n".join(stripper.parts)
                    break
    else:
        if msg.get_content_type().startswith("text/"):
            body = msg.get_content()

    text = "\n".join(header_lines) + ("\n\n" + body if body else "")
    return ExtractedDoc(
        text=text,
        metadata={"extractor": "eml", "headers": headers_meta},
    )


# ---------------------------------------------------------------------------
# CSV / TSV
# ---------------------------------------------------------------------------


def extract_csv(path: Path) -> ExtractedDoc:
    delim = "\t" if path.suffix.lower() == ".tsv" else ","
    rows: List[str] = []
    with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f, delimiter=delim)
        for row in reader:
            rows.append(" | ".join(row))
    return ExtractedDoc(
        text="\n".join(rows),
        metadata={"extractor": "csv", "row_count": len(rows)},
    )


# ---------------------------------------------------------------------------
# JSON — pretty-print, sorted keys (stable across re-ingests)
# ---------------------------------------------------------------------------


def extract_json(path: Path) -> ExtractedDoc:
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        text = json.dumps(json.loads(raw), indent=2, sort_keys=True, default=str)
    except json.JSONDecodeError:
        text = raw
    return ExtractedDoc(text=text, metadata={"extractor": "json"})


# ---------------------------------------------------------------------------
# Image / vision fallback — stub behind SOLOMON_ALLOW_VISION_API
# ---------------------------------------------------------------------------


def extract_text_via_sonnet(path: Path) -> str:  # noqa: ARG001
    """Stub for the multimodal-vision fallback.

    REPORT-CORPUS.md §4.6 documents this as the lone remote dependency.
    When ``SOLOMON_ALLOW_VISION_API=1`` is set, callers can override this
    by monkeypatching, or a future implementation can wire it up to the
    LLM client's vision tier. For now it always raises so scanned PDFs
    and images get parked rather than silently dropped.
    """
    if os.getenv("SOLOMON_ALLOW_VISION_API", "0") != "1":
        raise UnsupportedFileType(
            "vision fallback disabled (set SOLOMON_ALLOW_VISION_API=1 to enable)"
        )
    raise UnsupportedFileType(
        "vision fallback not yet implemented (stub). "
        "Configure a vision-capable model via the LLM client first."
    )


def extract_image(path: Path) -> ExtractedDoc:
    # HEIC requires pillow-heif at runtime. We accept it as part of IMAGE_EXT
    # so the routing layer doesn't have to special-case it.
    if path.suffix.lower() == ".heic":
        try:
            import pillow_heif  # type: ignore  # noqa: F401
        except ImportError as e:
            raise UnsupportedFileType("pillow-heif not installed for .heic") from e
    text = extract_text_via_sonnet(path)
    return ExtractedDoc(
        text=text,
        metadata={"extractor": "image", "via": "vision", "format": path.suffix.lower()},
    )


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------


_DISPATCH: List[Tuple[set, Any]] = [
    (PLAIN_TEXT_EXT, extract_txt_md),
    (RTF_EXT, extract_rtf),
    (PDF_EXT, extract_pdf),
    (DOCX_EXT, extract_docx),
    (PPTX_EXT, extract_pptx),
    (XLSX_EXT, extract_xlsx),
    (HTML_EXT, extract_html),
    (EMAIL_EXT, extract_eml),
    (CSV_EXT, extract_csv),
    (JSON_EXT, extract_json),
    (IMAGE_EXT, extract_image),
]


def extract(path: Path) -> ExtractedDoc:
    """Dispatch on suffix and return an ExtractedDoc.

    Raises
    ------
    UnsupportedFileType
        - Unknown extension
        - Required library missing
        - Extractor returned empty text
    """
    suffix = path.suffix.lower()
    for exts, fn in _DISPATCH:
        if suffix in exts:
            doc = fn(path)
            if not doc.text or not doc.text.strip():
                raise UnsupportedFileType(f"{suffix} extracted empty text")
            return doc
    raise UnsupportedFileType(f"{suffix} not recognized")
